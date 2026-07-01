"""
gerar_ppt.py — Mercury · Gerador de Propostas
Gera PPT simples com capa, resumo e tabela de fundos.
"""
from io import BytesIO
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

AZUL   = RGBColor(0x0a, 0x27, 0x44)
CINZA  = RGBColor(0xf5, 0xf5, 0xf2)
BRANCO = RGBColor(0xff, 0xff, 0xff)
VERDE  = RGBColor(0x1b, 0x5e, 0x20)

def _add_textbox(slide, text, left, top, width, height,
                 bold=False, size=12, color=None, align=PP_ALIGN.LEFT, bg=None):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    if bg:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg
    tf = txBox.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold  = bold
    run.font.size  = Pt(size)
    run.font.color.rgb = color or RGBColor(0x33,0x33,0x33)
    return txBox

def _fmt_pct(v):  return f"{v*100:.1f}%" if v is not None else "—"
def _fmt_brl(v):
    if v is None: return "—"
    return "R$ {:,.0f}".format(v).replace(",","X").replace(".",",").replace("X",".")

def gerar_ppt_proposta(cliente, valor, perfil, aloc_cl, aloc_fd, metricas):
    from fundos_config_v2 import CLASSES, fundos_por_classe

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank

    # ── SLIDE 1: CAPA ─────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    bg = s1.background.fill
    bg.solid()
    bg.fore_color.rgb = AZUL

    _add_textbox(s1, "MERCURY WEALTH MANAGEMENT",
                 0.5, 1.5, 12, 0.6, bold=True, size=14, color=BRANCO)
    _add_textbox(s1, "PROPOSTA DE CARTEIRA",
                 0.5, 2.2, 12, 1.0, bold=True, size=32, color=BRANCO)
    _add_textbox(s1, cliente.upper(),
                 0.5, 3.4, 12, 0.6, bold=True, size=18, color=RGBColor(0xcc,0xdd,0xee))
    _add_textbox(s1, f"Perfil: {perfil}   |   Valor: {_fmt_brl(valor)}   |   {date.today().strftime('%d/%m/%Y')}",
                 0.5, 4.2, 12, 0.5, size=12, color=RGBColor(0xaa,0xbb,0xcc))

    # ── SLIDE 2: RESUMO POR CLASSE ────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    _add_textbox(s2, "ALOCAÇÃO POR CLASSE DE ATIVO",
                 0.5, 0.3, 12, 0.5, bold=True, size=16, color=AZUL)

    headers = ["Classe de Ativo", "% Carteira", "Valor (R$)"]
    col_w   = [4.5, 1.5, 2.0]
    col_x   = [0.5, 5.2, 6.9]
    row_h   = 0.35
    y_start = 1.0

    # header row
    for i, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
        txb = s2.shapes.add_textbox(Inches(x), Inches(y_start), Inches(w), Inches(row_h))
        txb.fill.solid(); txb.fill.fore_color.rgb = AZUL
        p = txb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if i>0 else PP_ALIGN.LEFT
        r = p.add_run(); r.text = h
        r.font.bold=True; r.font.size=Pt(11); r.font.color.rgb=BRANCO

    y = y_start + row_h + 0.05
    for idx, c in enumerate(CLASSES):
        pct = aloc_cl.get(c, 0)
        if pct == 0: continue
        bg_c = RGBColor(0xf8,0xf8,0xf6) if idx%2==0 else BRANCO
        vals = [c, _fmt_pct(pct), _fmt_brl(pct*valor)]
        for i,(v,x,w) in enumerate(zip(vals, col_x, col_w)):
            txb = s2.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(row_h))
            txb.fill.solid(); txb.fill.fore_color.rgb = bg_c
            p = txb.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if i>0 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = v
            r.font.size=Pt(11); r.font.color.rgb=RGBColor(0x22,0x22,0x22)
        y += row_h + 0.02

    # total
    txb = s2.shapes.add_textbox(Inches(col_x[0]), Inches(y), Inches(3.5), Inches(row_h))
    txb.fill.solid(); txb.fill.fore_color.rgb = AZUL
    p = txb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "TOTAL"
    r.font.bold=True; r.font.size=Pt(11); r.font.color.rgb=BRANCO
    txb2 = s2.shapes.add_textbox(Inches(col_x[1]), Inches(y), Inches(1.5), Inches(row_h))
    txb2.fill.solid(); txb2.fill.fore_color.rgb = AZUL
    p2 = txb2.text_frame.paragraphs[0]; p2.alignment=PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "100,0%"
    r2.font.bold=True; r2.font.size=Pt(11); r2.font.color.rgb=BRANCO

    # ── SLIDE 3: FUNDOS ───────────────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    _add_textbox(s3, "FUNDOS DA CARTEIRA",
                 0.5, 0.3, 12, 0.5, bold=True, size=16, color=AZUL)

    headers3 = ["Fundo", "Classe", "% Total", "Valor", "Ret 6m", "Ret 12m", "Ret 24m", "YTD", "Vol 12m"]
    col3_w   = [4.0, 1.8, 0.9, 1.3, 0.8, 0.85, 0.85, 0.8, 0.85]
    col3_x   = [0.3]
    for w in col3_w[:-1]: col3_x.append(col3_x[-1]+w+0.02)

    y = 1.0
    for i,(h,x,w) in enumerate(zip(headers3, col3_x, col3_w)):
        txb = s3.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
        txb.fill.solid(); txb.fill.fore_color.rgb = AZUL
        p = txb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if i>0 else PP_ALIGN.LEFT
        r = p.add_run(); r.text = h
        r.font.bold=True; r.font.size=Pt(9); r.font.color.rgb=BRANCO

    y += 0.32
    row_h3 = 0.32
    for idx, c in enumerate(CLASSES):
        pct_c = aloc_cl.get(c, 0)
        if pct_c == 0: continue
        val_c = pct_c * valor
        for f in fundos_por_classe(c):
            nq   = f["nome_quantum"]
            p_fd_total = aloc_fd.get(nq, 0) / 100
            val_f = p_fd_total * valor
            if val_f == 0: continue
            m = metricas.get(nq, {})
            bg_r = RGBColor(0xf5,0xf5,0xf2) if idx%2==0 else BRANCO
            row_vals = [
                nq[:52], c,
                _fmt_pct(p_fd_total),
                _fmt_brl(val_f),
                _fmt_pct(m.get("ret_6m")),
                _fmt_pct(m.get("ret_12m")),
                _fmt_pct(m.get("ret_24m")),
                _fmt_pct(m.get("ytd")),
                _fmt_pct(m.get("vol_12m")),
            ]
            if y > 7.0: break
            for i,(v,x,w) in enumerate(zip(row_vals, col3_x, col3_w)):
                txb = s3.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(row_h3))
                txb.fill.solid(); txb.fill.fore_color.rgb = bg_r
                p = txb.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER if i>0 else PP_ALIGN.LEFT
                r = p.add_run(); r.text = v
                r.font.size=Pt(9); r.font.color.rgb=RGBColor(0x22,0x22,0x22)
            y += row_h3 + 0.01

    out = BytesIO()
    prs.save(out)
    return out.getvalue()
